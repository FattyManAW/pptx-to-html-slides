"""pptx_extractor.py — PPTX → JSON spec for gen_slides"""
import json, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt

def extract(pptx_path, section_names=None):
    """Extract PPTX content into structured JSON spec."""
    prs = Presentation(pptx_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        blocks = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        blocks.append(t)
            if shape.shape_type is not None and hasattr(shape, 'image'):
                try:
                    img = shape.image
                    if img.content:
                        ext = img.content_type.split('/')[-1]
                        if ext in ('png','jpg','jpeg','gif','webp'):
                            blocks.append(f"[IMAGE:{shape.name}.{ext}]")
                except Exception:
                    pass
        # Derive title from first non-empty block
        title = blocks[0] if blocks else ""
        body = " ".join(blocks[1:]) if len(blocks) > 1 else ""
        slides.append({"num": i+1, "title": title, "body": body, "blocks": blocks})
    # Section inference: every N slides or at explicit section breaks
    if section_names:
        n = len(section_names)
        per = max(1, len(slides) // n)
        for idx, sname in enumerate(section_names):
            pos = min(idx * per, len(slides) - 1)
            slides[pos]["_section"] = sname
    return {"total": len(slides), "slides": slides, "source": os.path.basename(pptx_path)}

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 pptx_extractor.py <pptx_path> [section_names_json]")
        sys.exit(1)
    pptx = sys.argv[1]
    secs = json.loads(sys.argv[2]) if len(sys.argv) > 2 else None
    spec = extract(pptx, secs)
    print(json.dumps(spec, ensure_ascii=False, indent=2))
