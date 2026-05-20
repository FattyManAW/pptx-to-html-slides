"""
extractor.py — PPTX → Structured JSON (gen_slides v4 module)

提供：
- extract(pptx_path) → structured slide spec (v4 format)
- classify_slide(blocks, slide_num) → slide type
- is_placeholder(text) → boolean
"""

import os
import sys
import re
import hashlib
from .themes import (
    PLACEHOLDER_PATTERNS,
    SECTION_KEYWORDS,
    COMPARISON_KEYWORDS,
)


def is_placeholder(text):
    """過濾 PPT 佔位符。"""
    for pat in PLACEHOLDER_PATTERNS:
        if re.search(pat, text):
            return True
    return False


def classify_slide(blocks, slide_num=1):
    """
    投影片分類 → "cover" | "section" | "thanks" | "twocol" | "items" | "content"
    """
    clean = [b.strip() for b in blocks if b.strip() and not is_placeholder(b)]
    if not clean:
        return "content"

    text = " ".join(clean)

    if slide_num == 1 and len(clean) <= 4:
        if not any(kw in text for kw in SECTION_KEYWORDS):
            return "cover"

    if any(kw in text for kw in ["Thank You", "谢谢", "謝謝", "感謝", "Thank you"]):
        return "thanks"

    if len(clean) <= 3:
        for kw in SECTION_KEYWORDS:
            if kw in text and len(text) <= 50:
                return "section"
        if re.match(r"^\d{2}\s", text) and len(text) <= 30:
            return "section"

    left_keywords = [p[0] for p in COMPARISON_KEYWORDS]
    right_keywords = [p[1] for p in COMPARISON_KEYWORDS]
    lcount = sum(1 for kw in left_keywords if kw in text)
    rcount = sum(1 for kw in right_keywords if kw in text)
    if lcount >= 1 and rcount >= 1:
        return "twocol"

    short_blocks = [b for b in clean if len(b) <= 30 and not re.match(r"^\d{1,2}$", b)]
    num_only = sum(1 for b in clean if re.match(r"^\d{1,2}$", b))
    if num_only > len(clean) * 0.4:
        return "content"
    if len(short_blocks) >= 5:
        return "items"

    return "content"


def extract(pptx_path, img_dir=None):
    """
    PPTX → 結構化 JSON（PLAYBOOK 2.5 格式）。
    """
    try:
        from pptx import Presentation
    except ImportError:
        print("❌ 需要 python-pptx: pip3 install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(pptx_path)
    slides_spec = []

    for i, slide in enumerate(prs.slides):
        blocks = []
        imgs = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t and not is_placeholder(t):
                        blocks.append(t)

            if hasattr(shape, "image"):
                try:
                    img = shape.image
                    if img.blob and img.content_type:
                        ext = img.content_type.split("/")[-1]
                        if ext in ("png", "jpg", "jpeg", "gif", "webp"):
                            if img_dir:
                                os.makedirs(img_dir, exist_ok=True)
                                name = shape.name.replace(" ", "_").replace("/", "_")
                                h = hashlib.md5(img.blob).hexdigest()[:8]
                                fname = f"slide{i+1}_{name}_{h}.{ext}"
                                fpath = os.path.join(img_dir, fname)
                                with open(fpath, "wb") as f:
                                    f.write(img.blob)
                                imgs.append(fname)
                            else:
                                imgs.append(f"{shape.name}.{ext}")
                except Exception:
                    pass

        stype = classify_slide(blocks, slide_num=i + 1)
        spec = {"num": i + 1, "type": stype}

        if stype == "cover":
            spec["t"] = blocks[0] if blocks else ""
            spec["st"] = blocks[1] if len(blocks) > 1 else ""
            spec["blocks"] = blocks
        elif stype == "section":
            spec["n"] = f"{i+1:02d}"
            spec["t"] = blocks[0] if blocks else ""
            spec["blocks"] = blocks
        elif stype == "thanks":
            spec["t"] = blocks[0] if blocks else ""
            spec["st"] = blocks[1] if len(blocks) > 1 else ""
            spec["blocks"] = blocks
        elif stype == "twocol":
            spec["t"] = blocks[0] if blocks else ""
            half = max(1, len(blocks[1:]) // 2) if len(blocks) > 1 else 0
            spec["twocol"] = [
                {"side": "left", "ls": blocks[1:1+half] if half > 0 else []},
                {"side": "right", "ls": blocks[1+half:] if half > 0 else []},
            ]
            spec["blocks"] = blocks
        elif stype == "items":
            spec["t"] = blocks[0] if blocks else ""
            card_blocks = blocks[1:]
            cards = []
            j = 0
            while j < len(card_blocks):
                card = {"h": card_blocks[j]}
                details = []
                j += 1
                while j < len(card_blocks) and len(card_blocks[j]) > 30:
                    details.append(card_blocks[j])
                    j += 1
                if details:
                    card["ls"] = details
                cards.append(card)
            spec["items"] = cards
            spec["blocks"] = blocks
        else:  # content
            spec["t"] = blocks[0] if blocks else ""
            if len(blocks) > 1:
                spec["b"] = "\n".join(blocks[1:2])
            spec["blocks"] = blocks

        if imgs:
            spec["imgs"] = imgs
        slides_spec.append(spec)

    return {
        "version": "v4",
        "source": os.path.basename(pptx_path),
        "total": len(slides_spec),
        "slides": slides_spec,
    }


# 向後相容：extract_pptx 別名
extract_pptx = extract